"""
Indiana Basketball — 2026 Portal Production Report by Position
--------------------------------------------------------------
Clean light-background design, IU crimson.
Reads from data/d1_master_2026.csv (copied from OSUPortal; 2025-26 season data).
Groups the 2026 transfer-portal class into Guards / Wings / Bigs and shows the
top 15 per group ranked by current-season BPR (production). No projections.

Run:
  python deck/build_position_deck.py
"""

import os, sys, time, requests
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

# ── Color palette (IU crimson — constant names kept from the OSU template) ────
SCARLET    = RGBColor(0x99, 0x00, 0x00)   # IU crimson #990000
SCARLET_DK = RGBColor(0x66, 0x00, 0x00)
SCARLET_BG = RGBColor(0xF5, 0xEA, 0xEA)   # very light crimson tint
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PANEL      = RGBColor(0xF7, 0xF8, 0xFA)
PANEL2     = RGBColor(0xEE, 0xEF, 0xF1)
BORDER     = RGBColor(0xD8, 0xD8, 0xD8)
TEXT       = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MED   = RGBColor(0x50, 0x50, 0x50)
TEXT_LIGHT = RGBColor(0x96, 0x96, 0x96)
GREEN      = RGBColor(0x1A, 0x80, 0x3E)
GREEN_BAR  = RGBColor(0x27, 0xAE, 0x60)
GREEN_BG   = RGBColor(0xE4, 0xF6, 0xEB)
RED        = RGBColor(0xBB, 0x33, 0x26)
RED_BAR    = RGBColor(0xE7, 0x4C, 0x3C)
RED_BG     = RGBColor(0xFC, 0xEB, 0xEA)
BLUE_BAR   = RGBColor(0x4A, 0x90, 0xC0)
GOLD       = RGBColor(0x5A, 0x5A, 0x5A)   # OSU gray (replaces gold)
GOLD_LT    = RGBColor(0xA8, 0xA8, 0xA8)   # lighter gray for dark backgrounds
GOLD_BG    = RGBColor(0xF0, 0xF0, 0xF0)   # light gray background
OSU_GRAY    = RGBColor(0x41, 0x41, 0x41)  # OSU official gray
OSU_GRAY_DK = RGBColor(0x2C, 0x2C, 0x2C)

P6        = {'ACC', 'B10', 'B12', 'P12', 'SEC', 'BE'}
TOP_N_PER_GROUP = 15
MAX_GAMES = 40     # sanity ceiling — no D1 team plays more games than this in a season

# Position groups (BartTorvik roles → PG / Combo / Wing / PF / C)
POSITION_GROUPS = [
    ("PG",    {'Pure PG', 'Scoring PG'}),
    ("Combo", {'Combo G'}),
    ("Wing",  {'Wing G', 'Wing F'}),
    ("PF",    {'Stretch 4', 'PF/C'}),
    ("C",     {'C'}),
]
GROUP_ROMAN = {"PG": "I", "Combo": "II", "Wing": "III", "PF": "IV", "C": "V"}
GROUP_ROLES_DESC = {
    "PG":    "Pure PG  ·  Scoring PG",
    "Combo": "Combo G",
    "Wing":  "Wing G  ·  Wing F",
    "PF":    "Stretch 4  ·  PF/C",
    "C":     "Center",
}
ROLE_ORDER = {'Scoring PG': 0, 'Pure PG': 1, 'Combo G': 2, 'Wing G': 3,
              'Wing F': 4, 'Stretch 4': 5, 'PF/C': 6, 'C': 7}

# Players to exclude (case-insensitive match on ESPN name). None for this report.
EXCLUDE_PLAYERS = set()

# Teams to exclude — players who entered the portal FROM these schools are not
# candidates in this report (e.g. Indiana's own outgoing transfers, such as
# Tayton Conerway, shouldn't appear as "targets" in Indiana's own portal report).
EXCLUDE_FROM_TEAMS = {"Indiana Hoosiers"}

# ── Contact info shown on title slide (leave blank to omit) ───────────────────
CONTACT_NAME  = "Armeen Ghoorkhanian"
CONTACT_EMAIL = "ghoorkhanian.5@osu.edu"
CONTACT_PHONE = "614-578-9275"

# ── Shared table layout ────────────────────────────────────────────────────────
TABLE_HDRS = ["#","Player","School","Conf","Role","Class","Ht",
              "Pts/40","Reb/40","Ast/40","eFG%","USG%","3PT%",
              "OBPR","DBPR","BPR"]
TABLE_WS   = [0.44, 1.78, 2.60, 0.64, 0.84, 0.52, 0.56,
              0.60, 0.60, 0.60, 0.54, 0.54, 0.52,
              0.52, 0.52, 0.92]
# Total ≈ 12.74" — fits in 13.33" slide with 0.12 margins

ROOT      = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA_DIR  = os.path.join(ROOT, "data")
OUTPUT    = os.path.join(ROOT, "output")
IMG_DIR   = os.path.join(DATA_DIR, "headshots")
os.makedirs(IMG_DIR, exist_ok=True)
os.makedirs(OUTPUT, exist_ok=True)
MASTER_CSV = os.path.join(DATA_DIR, "d1_master_2026.csv")

CONF_LABELS = {
    'WCC':'West Coast','MWC':'Mountain West','A10':'Atlantic 10',
    'Amer':'American','MVC':'Missouri Valley','WAC':'Western Athletic',
    'Ivy':'Ivy League','CUSA':'Conf. USA','Slnd':'Southland',
    'MAC':'Mid-American','CAA':'Colonial','Horz':'Horizon',
    'Sum':'Summit League','BSky':'Big Sky','SC':'Southern',
    'SB':'Sun Belt','OVC':'Ohio Valley','NEC':'Northeast',
    'Pat':'Patriot','MEAC':'MEAC','SWAC':'SWAC',
    'AE':'America East','ASun':'ASUN','BW':'Big West',
    'BSth':'Big South','ACC':'ACC','B10':'Big Ten',
    'B12':'Big 12','SEC':'SEC','BE':'Big East','P12':'Pac-12',
}
ROLE_DISPLAY = {
    'Scoring PG': 'Scoring PG', 'Pure PG': 'Pure PG',
    'Combo G': 'Combo G',       'Wing G': 'Wing/G',
    'Wing F': 'Wing/F',         'Stretch 4': 'Stretch 4',
    'PF/C': 'PF/C',             'C': 'Center',
}
ROLE_DESC = {
    'Scoring PG': 'scoring point guard', 'Pure PG': 'lead playmaker',
    'Combo G': 'combo guard',            'Wing G': 'wing/guard',
    'Wing F': 'wing/forward',            'Stretch 4': 'stretch four',
    'PF/C': 'power forward/center',      'C': 'center',
}
# stat key → (label, higher_is_better)
STAT_META = {
    'pts_40':    ('Points / 40',    True),
    'reb_40':    ('Rebounds / 40',  True),
    'ast_40':    ('Assists / 40',   True),
    'stl_40':    ('Steals / 40',    True),
    'tov_40':    ('Turnovers / 40', False),
    'mpg':       ('Min Per Game',   True),
    'efg_pct':   ('eFG %',          True),
    'usg_pct':   ('Usage %',        True),
    'three_pct': ('3-Point %',      True),
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def R(slide, l, t, w, h, clr, border_clr=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = clr
    if border_clr:
        s.line.color.rgb = border_clr
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def T(slide, text, l, t, w, h, sz=10, bold=False, clr=TEXT,
      align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text); run.font.size = Pt(sz)
    run.font.bold   = bold
    run.font.color.rgb = clr
    run.font.italic = italic
    run.font.name   = "Calibri"
    return box


def pct_bar(slide, pct_val, l, t, w, h=0.055, higher_better=True):
    R(slide, l, t, w, h, BORDER)
    if pct_val is None: return
    eff = pct_val if higher_better else (100 - pct_val)
    fill_w = max(w * eff / 100, 0.004)
    clr = GREEN_BAR if eff >= 80 else (RED_BAR if eff <= 20 else BLUE_BAR)
    R(slide, l, t, min(fill_w, w), h, clr)


def fmt(v, dec=1, suf=""):
    try:
        f = float(v)
        if f != f:  # NaN
            return f"{0:.{dec}f}{suf}"
        return f"{f:.{dec}f}{suf}"
    except: return "—"


def ordinal(n):
    """Return '1st', '2nd', '83rd', '42nd', etc."""
    try:
        n = int(n)
        if 11 <= (n % 100) <= 13:
            return f"{n}th"
        return f"{n}" + {1:'st', 2:'nd', 3:'rd'}.get(n % 10, 'th')
    except:
        return "—"


def bpr_clr(v, dark_bg=False):
    try:
        b = float(v)
        if dark_bg:
            if b >= 7:  return RGBColor(0xFF, 0xCC, 0xCC)
            if b >= 5:  return GOLD_LT
            if b >= 3:  return WHITE
            return RGBColor(0xFF, 0xAA, 0xAA)
        else:
            if b >= 7:  return SCARLET
            if b >= 5:  return GOLD
            if b >= 3:  return TEXT
            return RED
    except: return TEXT_LIGHT


def val_color(pct_val, hib=True):
    if pct_val is None: return TEXT_LIGHT
    eff = pct_val if hib else (100 - pct_val)
    if eff >= 80: return GREEN
    if eff <= 20: return RED
    return TEXT


def years_remaining(class_yr):
    return {'Fr':'FR','So':'SO','Jr':'JR','Sr':'SR','Gr':'GR'}.get(
        str(class_yr).strip(), '—')


def elig_clr(class_yr):
    return GOLD if str(class_yr).strip() == 'Sr' else TEXT


def _section_tag(slide, group_name):
    """Minimal section indicator tucked into the top-right of the crimson header band."""
    label = f"SECTION {GROUP_ROMAN.get(group_name, '')}  ·  {group_name.upper()}"
    T(slide, label, 10.60, 0.10, 2.55, 0.22,
      sz=7, bold=True, clr=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.RIGHT)


# ── Percentiles ───────────────────────────────────────────────────────────────

def build_pct_arrays(d1_df):
    base = d1_df[pd.to_numeric(d1_df["games"], errors="coerce").fillna(0) >= 10]
    out  = {}
    for col in ["pts_40","reb_40","ast_40","stl_40","blk_40","tov_40",
                "fg_pct","three_pct","ft_pct","efg_pct","usg_pct","mpg",
                "obpr","dbpr","bpr"]:
        vals = pd.to_numeric(base.get(col, pd.Series(dtype=float)),
                             errors="coerce").dropna().values
        out[col] = np.sort(vals) if len(vals) else np.array([])
    return out


def get_pct(val, arr):
    if arr is None or len(arr) == 0: return None
    try:
        v = float(val)
        return round(np.searchsorted(arr, v, side="right") / len(arr) * 100, 1)
    except: return None


# ── Scout context ─────────────────────────────────────────────────────────────

def generate_scout_context(row):
    """Returns 2 (label, value, color) tuples for the Scout Context panel."""
    rows_out = []

    # ── 1. Team Quality ────────────────────────────────────────────────────────
    # Thresholds calibrated from 2025-26 BartTorvik data (363 D1 teams):
    # Top 25 ≈ net ≥ 22.4 | Top 50 ≈ net ≥ 15.4 | Top 100 ≈ net ≥ 7.1
    # Top 125 ≈ net ≥ 3.6 | Top 150 ≈ net ≥ 1.0 | Top 200 ≈ net ≥ -2.4
    team_net = row.get('team_net')
    if pd.notna(team_net):
        net_val = float(team_net)
        if net_val >= 22:
            tier = "Top 25 nationally"
        elif net_val >= 15:
            tier = "Top 50 nationally"
        elif net_val >= 7:
            tier = "Top 100 nationally"
        elif net_val >= 3.5:
            tier = "Top 125 nationally"
        elif net_val >= 1:
            tier = "Top 150 nationally"
        elif net_val >= -2.5:
            tier = "Top 200 nationally"
        else:
            tier = "Lower-tier program"
        sign = "+" if net_val >= 0 else ""
        net_str = f"AdjNet {sign}{net_val:.1f}  \u00b7  {tier}"
        net_clr = GREEN if net_val >= 15 else GOLD if net_val >= 7 else TEXT_MED
    else:
        net_str = "\u2014"
        net_clr = TEXT_MED
    rows_out.append(("Team Quality", net_str, net_clr))

    # ── 2. Eligibility ─────────────────────────────────────────────────────────
    # Data is from 2025-26; advance by one year for 2026-27 projection.
    # Fr→So, So→Jr, Jr→Sr, Sr→Gr, Gr→5th year
    cy = str(row.get('class_year') or '').strip()
    cy_lower = cy.lower()
    if 'grad' in cy_lower or cy_lower in ('g', 'gr'):
        elig = "5th Year / Grad  \u00b7  1 year remaining"
        elig_clr = GOLD
    elif 'sr' in cy_lower or 'senior' in cy_lower:
        elig = "Graduate Transfer  \u00b7  1 year remaining"
        elig_clr = GOLD
    elif 'jr' in cy_lower or 'junior' in cy_lower:
        elig = "Senior  \u00b7  1 year remaining"
        elig_clr = GOLD
    elif 'so' in cy_lower or 'sophomore' in cy_lower:
        elig = "Junior  \u00b7  2 years remaining"
        elig_clr = TEXT_MED
    elif 'fr' in cy_lower or 'freshman' in cy_lower:
        elig = "Sophomore  \u00b7  3 years remaining"
        elig_clr = TEXT_MED
    else:
        elig = cy if cy and cy.lower() not in ('nan', 'none', '') else "\u2014"
        elig_clr = TEXT_MED
    rows_out.append(("Eligibility", elig, elig_clr))

    return rows_out


# ── Headshots ─────────────────────────────────────────────────────────────────

def fetch_headshot(espn_id, headshot_url=None):
    if not espn_id or str(espn_id) in ("nan","None",""): return None
    path = os.path.join(IMG_DIR, f"{espn_id}.png")
    if os.path.exists(path): return path
    clean = headshot_url if headshot_url and str(headshot_url) not in ("nan","None","") else None
    url = clean or f"https://a.espncdn.com/i/headshots/mens-college-basketball/players/full/{espn_id}.png"
    try:
        r = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        if r.ok and len(r.content) > 4000:
            with open(path, "wb") as f: f.write(r.content)
            time.sleep(0.25)
            return path
    except: pass
    return None


def add_photo(slide, path, l, t, w, h):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
            return
        except: pass
    R(slide, l, t, w, h, PANEL2, BORDER)
    T(slide, "No photo", l, t + h/2 - 0.14, w, 0.28,
      sz=8, clr=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ── Title slide ───────────────────────────────────────────────────────────────

def slide_title(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)

    # Scarlet left accent bar
    R(s, 0, 0, 0.08, 7.5, SCARLET_DK)

    # Large title — TRANSFER PORTAL stacked
    T(s, "TRANSFER", 0.45, 0.95, 12, 1.10, sz=72, bold=True, clr=TEXT)
    T(s, "PORTAL",   0.45, 1.99, 12, 1.10, sz=72, bold=True, clr=TEXT)

    # Scarlet rule under title
    R(s, 0.45, 3.21, 1.20, 0.055, SCARLET)

    # Subtitle
    T(s, "Production Report by Position", 0.45, 3.39, 10, 0.54, sz=28, clr=TEXT_MED)
    T(s, "Top 15 PG  ·  Top 15 Combo  ·  Top 15 Wing  ·  Top 15 PF  ·  Top 15 C  ·  Ranked by Career Production (BPR)",
      0.45, 4.05, 11, 0.32, sz=13, clr=TEXT_LIGHT)

    # Year — light gray, far right, top
    T(s, "2026 Portal Class", 9.0, 0.82, 4.15, 0.36,
      sz=16, bold=True, clr=BORDER, align=PP_ALIGN.RIGHT)

    # Footer strip
    R(s, 0, 6.97, 13.33, 0.53, PANEL)
    R(s, 0, 6.97, 13.33, 0.008, BORDER)
    T(s, "2025\u201326 Season Data",
      0.45, 7.10, 10, 0.28, sz=9, clr=TEXT_LIGHT)

    # Contact info — right-aligned in footer (only rendered if at least one field set)
    contact_parts = [p for p in [CONTACT_NAME, CONTACT_EMAIL, CONTACT_PHONE] if p.strip()]
    if contact_parts:
        T(s, "  \u00b7  ".join(contact_parts),
          3.0, 7.10, 10.10, 0.28, sz=9, clr=TEXT_MED, align=PP_ALIGN.RIGHT)


# ── Explainer slides ──────────────────────────────────────────────────────────

def _lerp_color(c1, c2, t):
    """Linearly interpolate between two (r, g, b) tuples (t in [0, 1])."""
    return RGBColor(*(int(c1[k] + (c2[k] - c1[k]) * t) for k in range(3)))


def slide_explainer_structure(prs):
    """Slide 2A — This Report Has Five Sections (PG / Combo / Wing / PF / C)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)

    n = len(POSITION_GROUPS)
    PANEL_W = 13.33 / n

    # Gradient from IU crimson to gray across however many panels there are.
    bg_ends       = ((0x99, 0x00, 0x00), (0x2C, 0x2C, 0x2C))
    numeral_ends  = ((0xB3, 0x22, 0x22), (0x55, 0x55, 0x55))
    accent_ends   = ((0xFF, 0xAA, 0xAA), (0xAA, 0xAA, 0xAA))
    sub_ends      = ((0xFF, 0xCC, 0xCC), (0xBB, 0xBB, 0xBB))

    def _ramp(ends):
        return [_lerp_color(ends[0], ends[1], i / (n - 1) if n > 1 else 0)
                for i in range(n)]

    panel_bgs   = _ramp(bg_ends)
    numeral_clr = _ramp(numeral_ends)
    accent_clr  = _ramp(accent_ends)
    sub_clr     = _ramp(sub_ends)

    # Scale the decorative numeral down as columns get narrower (tuned for n=3).
    numeral_sz = max(60, round(150 * (PANEL_W - 0.45) / (13.33 / 3 - 0.45)))

    for i, (group_name, _) in enumerate(POSITION_GROUPS):
        px = i * PANEL_W
        R(s, px, 0, PANEL_W, 7.5, panel_bgs[i])

        # Decorative large section numeral
        T(s, GROUP_ROMAN[group_name], px + 0.15, 0.40, PANEL_W - 0.45, 5.50,
          sz=numeral_sz, bold=True, clr=numeral_clr[i], align=PP_ALIGN.RIGHT)

        LX = px + 0.42
        LW = PANEL_W - 0.70
        T(s, f"SECTION  {GROUP_ROMAN[group_name]}", LX, 0.50, LW, 0.28,
          sz=10, bold=True, clr=accent_clr[i])
        T(s, group_name, LX, 0.88, LW, 0.90, sz=34, bold=True, clr=WHITE)

        R(s, LX, 1.90, min(1.30, LW), 0.04, accent_clr[i])

        T(s, "This past year's portal class,\nranked by career production",
          LX, 2.10, LW, 0.56, sz=10, clr=sub_clr[i])

        R(s, LX, 6.50, LW, 0.56, RGBColor(0x1E, 0x1E, 0x1E))
        T(s, "15 players", LX + 0.14, 6.55, LW - 0.24, 0.26,
          sz=13, bold=True, clr=WHITE)
        T(s, "Sorted by career\nproduction (BPR)", LX + 0.14, 6.80, LW - 0.24, 0.30,
          sz=7.5, clr=accent_clr[i])

    # Full-width bottom strip
    R(s, 0, 7.15, 13.33, 0.35, RGBColor(0x1E, 0x1E, 0x1E))
    T(s, "2026 TRANSFER PORTAL PRODUCTION REPORT",
      0.30, 7.19, 13.0, 0.26, sz=7.5, bold=True,
      clr=TEXT_LIGHT, align=PP_ALIGN.CENTER)


def slide_explainer_bpr(prs):
    """Slide 2B — What is BPR?"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)

    # Scarlet header strip
    R(s, 0, 0, 13.33, 0.60, SCARLET)
    T(s, "WHAT IS BPR?", 0.22, 0.10, 9, 0.44, sz=17, bold=True, clr=WHITE)
    T(s, "Bayesian Performance Rating — one number for total impact",
      8.0, 0.18, 5.15, 0.28, sz=9, clr=WHITE, align=PP_ALIGN.RIGHT, italic=True)

    BODY_Y = 0.72
    LEFT_W = 5.60
    RIGHT_W = 7.20
    RX = 0.20 + LEFT_W + 0.20

    # ── Left column: 3 bullets ────────────────────────────────────────────────
    LX = 0.20
    y = BODY_Y

    T(s, "BPR IN PLAIN LANGUAGE", LX, y, LEFT_W, 0.28,
      sz=9, bold=True, clr=SCARLET)
    y += 0.38

    bullets_bpr = [
        ("One number captures the full picture",
         "Instead of tracking 10 separate stats, BPR combines everything — offense, defense, efficiency — into a single score that measures how much a player contributes to winning."),
        ("0 = average D1 player. Higher is better.",
         "A score of 0 means exactly average across all of D1. Positive = above average, negative = below. The higher the number, the more impact they create."),
        ("Combines offense AND defense",
         "BPR is built from two components: Offensive BPR (scoring, assists, efficiency) and Defensive BPR (stops, steals, positioning). Both matter equally in the final score."),
    ]

    for header, body in bullets_bpr:
        R(s, LX + 0.04, y + 0.10, 0.10, 0.10, SCARLET)
        T(s, header, LX + 0.24, y, LEFT_W - 0.28, 0.26,
          sz=10, bold=True, clr=TEXT)
        T(s, body, LX + 0.24, y + 0.28, LEFT_W - 0.28, 0.50,
          sz=8.5, clr=TEXT_MED)
        y += 0.90

    # How the ranking works callout
    R(s, LX, y + 0.10, LEFT_W, 0.010, BORDER)
    y += 0.24
    T(s, "How this report ranks players:", LX, y, LEFT_W, 0.26,
      sz=9, bold=True, clr=TEXT)
    y += 0.32

    for box_clr, label, sub in [
        (SCARLET_BG, "Career BPR",  "Every portal player is ranked by their career production"),
        (PANEL2,     "By position",  "Top 15 Guards, top 15 Wings, top 15 Bigs"),
    ]:
        R(s, LX, y, LEFT_W, 0.56, box_clr, BORDER)
        T(s, label, LX + 0.16, y + 0.07, LEFT_W - 0.24, 0.22,
          sz=9, bold=True, clr=SCARLET)
        T(s, sub, LX + 0.16, y + 0.30, LEFT_W - 0.24, 0.20,
          sz=8, clr=TEXT_MED)
        y += 0.62

    # ── Right column: Rating scale table ──────────────────────────────────────
    y = BODY_Y
    T(s, "BPR RATING SCALE", RX, y, RIGHT_W, 0.28,
      sz=9, bold=True, clr=SCARLET)
    y += 0.38

    # Table header — two columns only (range + description/example combined)
    R(s, RX, y, RIGHT_W, 0.32, SCARLET_DK)
    T(s, "BPR Range", RX + 0.12, y + 0.06, 1.20, 0.22,
      sz=8, bold=True, clr=WHITE)
    T(s, "What it means  ·  reference examples", RX + 1.50, y + 0.06, 5.50, 0.22,
      sz=8, bold=True, clr=WHITE)
    y += 0.32

    scale_rows = [
            ("8.0 +",     "Elite / All-Conference — Game-changing impact",      "Cameron Boozer, Duke (14.5)  ·  Milan Momcilovic, Iowa St (8.5)", SCARLET,    WHITE),
            ("5.0 – 8.0", "Quality P6 Starter — Highly positive net impact",    "Aiden Sherrell, Alabama (7.2)  ·  Markus Burton, Notre Dame (5.6)", GOLD,       PANEL),
            ("2.0 – 5.0", "Starter/Quality Rotation Piece — Solid contributor at the P6 level", "Samet Yigitoglu, SMU (4.3)  ·  Bryce Lindsay, Villanova (3.6)", TEXT,       WHITE),
            ("0.0 – 2.0", "High Major Bench Piece / Mid Major Starter",   "Darren Harris, Duke (1.5)  ·  Jaeden Mustaf, Georgia Tech (1.3)", TEXT_LIGHT, PANEL),
            ("< 0.0",     "Project / Deep Bench — Below average D1 player",     "", TEXT_LIGHT, WHITE),
        ]

    for rng, desc, ex, txt_clr, row_bg in scale_rows:
        R(s, RX, y, RIGHT_W, 0.56, row_bg, BORDER)
        R(s, RX, y, 1.34, 0.56, row_bg)
        T(s, rng,  RX + 0.12, y + 0.16, 1.10, 0.26,
          sz=11, bold=True, clr=txt_clr)
        T(s, desc, RX + 1.50, y + 0.08, 5.50, 0.22,
          sz=8.5, clr=TEXT)
        T(s, ex,   RX + 1.50, y + 0.30, 5.50, 0.20,
          sz=7.5, clr=TEXT_LIGHT, italic=True)
        y += 0.56

    y += 0.20
    R(s, RX, y, RIGHT_W, 0.010, BORDER)
    y += 0.18

    T(s, "STAT COLORS ON PLAYER CARDS", RX, y, RIGHT_W, 0.26,
      sz=8.5, bold=True, clr=SCARLET)
    y += 0.32

    T(s, "Each stat bar shows where the player ranks vs all D1 players this season:",
      RX, y, RIGHT_W, 0.24, sz=8.5, clr=TEXT_MED)
    y += 0.32

    for bar_clr, label in [
        (GREEN_BAR, "Green — top 20% of D1"),
        (BLUE_BAR,  "Blue — middle of the pack"),
        (RED_BAR,   "Red — bottom 20% of D1"),
    ]:
        R(s, RX + 0.10, y + 0.06, 0.28, 0.14, bar_clr)
        T(s, label, RX + 0.48, y, RIGHT_W - 0.56, 0.26,
          sz=8.5, clr=TEXT_MED)
        y += 0.28


# ── Shared table renderer ─────────────────────────────────────────────────────

def _draw_targets_table(s, df, x0, y0, rh, pct_arrays):
    """Render the shared top-15 table used by each position-group table slide."""
    hdrs = TABLE_HDRS
    ws   = TABLE_WS

    STAT_PCT = {
        7:  ("pts_40",    True),
        8:  ("reb_40",    True),
        9:  ("ast_40",    True),
        10: ("efg_pct",   True),
        11: ("usg_pct",   True),
        12: ("three_pct", True),
    }

    # Header row
    R(s, x0, y0, sum(ws), rh, SCARLET_DK)
    x = x0
    for h, w in zip(hdrs, ws):
        T(s, h, x+0.04, y0+0.05, w-0.06, rh-0.02, sz=7, bold=True, clr=WHITE)
        x += w

    cur_y   = y0 + rh
    row_num = 0

    for _, row in df.iterrows():
        role = str(row.get("role","") or "")

        row_num += 1
        cls       = str(row.get("class_year","") or "")
        role_disp = ROLE_DISPLAY.get(role, role)[:12]

        obpr = float(row.get("obpr") or 0)
        dbpr = float(row.get("dbpr") or 0)
        bpr  = obpr + dbpr

        vals = [
            str(row_num),
            str(row.get("espn_name",""))[:22],
            str(row.get("espn_team",""))[:30],
            str(row.get("conf",""))[:8],
            role_disp,
            years_remaining(cls),
            str(row.get("height_display") or "—"),
            fmt(row.get("pts_40")),
            fmt(row.get("reb_40")),
            fmt(row.get("ast_40")),
            fmt(row.get("efg_pct"),1,"%"),
            fmt(row.get("usg_pct"),1,"%"),
            fmt(row.get("three_pct"),1,"%"),
            fmt(obpr, 2),
            fmt(dbpr, 2),
            fmt(bpr, 2),
        ]

        def _sc(col_idx, raw_val):
            if col_idx not in STAT_PCT:
                return TEXT
            key, hib = STAT_PCT[col_idx]
            pct = get_pct(raw_val, pct_arrays.get(key))
            return val_color(pct, hib)

        clrs = [
            TEXT_LIGHT,
            TEXT,
            TEXT_MED,
            TEXT_LIGHT,
            SCARLET,
            TEXT_LIGHT,
            TEXT_LIGHT,
            _sc(7,  row.get("pts_40")),
            _sc(8,  row.get("reb_40")),
            _sc(9,  row.get("ast_40")),
            _sc(10, row.get("efg_pct")),
            _sc(11, row.get("usg_pct")),
            _sc(12, row.get("three_pct")),
            TEXT_MED,
            TEXT_MED,
            TEXT_MED,
        ]

        bg = WHITE if row_num % 2 == 0 else PANEL
        R(s, x0, cur_y, sum(ws), rh, bg)
        R(s, x0, cur_y + rh - 0.004, sum(ws), 0.004, BORDER)
        x = x0
        for vi, (val, w, clr) in enumerate(zip(vals, ws, clrs)):
            T(s, val, x+0.04, cur_y+0.04, w-0.06, rh-0.04,
              sz=7.5, clr=clr, bold=(vi == len(vals)-1))
            x += w
        cur_y += rh


# ── Position-group table slide ────────────────────────────────────────────────

def slide_group_table(prs, group_name, df, pct_arrays):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)
    R(s, 0, 0, 13.33, 0.60, SCARLET)

    T(s, f"{group_name.upper()} — TOP {len(df)} BY PRODUCTION",
      0.20, 0.08, 9, 0.46, sz=17, bold=True, clr=WHITE)
    T(s, f"{GROUP_ROLES_DESC[group_name]}  ·  Sorted by career production (BPR)",
      6.60, 0.36, 6.55, 0.28, sz=8.5, clr=WHITE, align=PP_ALIGN.RIGHT, italic=True)
    _section_tag(s, group_name)

    _draw_targets_table(s, df.head(TOP_N_PER_GROUP), x0=0.12, y0=0.66,
                        rh=0.385, pct_arrays=pct_arrays)

    T(s, "Stat colors: green = top 20% D1  ·  red = bottom 20%",
      0.12, 7.22, 13.0, 0.24, sz=6.5, clr=TEXT_LIGHT, italic=True)


# ── Player slide ──────────────────────────────────────────────────────────────

def slide_player(prs, rank, row, pct_arrays, n_total, group_name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)

    # ── HEADER ────────────────────────────────────────────────────────────────
    HDR_H = 1.15
    R(s, 0, 0, 13.33, HDR_H, SCARLET)
    R(s, 0, 0, 0.14, HDR_H, SCARLET_DK)
    _section_tag(s, group_name)

    # Rank chip
    chip_sz = 14 if rank >= 10 else 16
    R(s, 0.16, 0.28, 0.62, 0.56, SCARLET_DK)
    T(s, f"#{rank}", 0.16, 0.38, 0.62, 0.32,
      sz=chip_sz, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

    # Headshot
    espn_id = str(row.get("espn_id",""))
    hs_url  = str(row.get("headshot_url",""))
    hs_path = fetch_headshot(espn_id, hs_url)
    add_photo(s, hs_path, 0.82, 0.18, 1.00, 0.80)

    # Player attributes
    player   = str(row.get("espn_name",""))
    team     = str(row.get("espn_team",""))
    conf     = CONF_LABELS.get(str(row.get("conf","")), str(row.get("conf","")))
    ht       = str(row.get("height_display","") or "")
    wt       = str(row.get("weight_display","") or "")
    cls      = str(row.get("class_year","") or "")
    role     = str(row.get("role","") or "")
    birth    = str(row.get("birthplace","") or "")
    cls_abbr = {'Fr':'FR','So':'SO','Jr':'JR','Sr':'SR','Gr':'GR'}.get(
        cls.strip(), cls[:2].upper() if cls.strip() not in ("nan","") else "—")

    T(s, player, 1.96, 0.06, 9.00, 0.48, sz=22, bold=True, clr=WHITE)

    sub1_parts = [p for p in [team, conf, ht, wt] if p and p not in ("nan","","—")]
    T(s, "  ·  ".join(sub1_parts), 1.96, 0.54, 9.00, 0.26, sz=9.5, clr=WHITE)

    role_disp = ROLE_DISPLAY.get(role, role)
    sub2_parts = []
    if role_disp and role_disp not in ("nan","","—"):
        sub2_parts.append(role_disp)
    if cls_abbr and cls_abbr not in ("—","nan",""):
        sub2_parts.append(cls_abbr)
    if birth and birth not in ("nan","","—"):
        sub2_parts.append(f"📍 {birth}")
    T(s, "  ·  ".join(sub2_parts), 1.96, 0.83, 10.50, 0.26, sz=9, clr=WHITE)

    # ── BODY ──────────────────────────────────────────────────────────────────
    BY = HDR_H + 0.07
    BH = 7.5 - BY - 0.14

    C1L, C1W = 0.12, 7.20
    C2L, C2W = 7.42, 5.79

    R(s, C1L, BY, C1W, BH, WHITE, BORDER)
    R(s, C2L, BY, C2W, BH, WHITE, BORDER)

    PAD = 0.15

    # ── COL 1: Stats ──────────────────────────────────────────────────────────
    sx = C1L + PAD
    sw = C1W - PAD*2
    sy = BY + PAD

    T(s, "2025–26 STATS", sx, sy, sw, 0.28, sz=9, bold=True, clr=SCARLET)
    T(s, "Per 40 minutes  ·  vs all D1 players (10+ games)",
      sx, sy+0.28, sw, 0.20, sz=7, clr=TEXT_LIGHT, italic=True)

    stats = [
        ('mpg',       True),
        ('pts_40',    True),
        ('reb_40',    True),
        ('ast_40',    True),
        ('stl_40',    True),
        ('tov_40',    False),
        ('efg_pct',   True),
        ('usg_pct',   True),
        ('three_pct', True),
    ]
    n_stats     = len(stats)
    available_h = BH - 0.56
    ROW_H       = min(available_h / (n_stats + 0.8), 0.58)
    st_y        = sy + 0.54

    for i, (key, hib) in enumerate(stats):
        if i == 1:
            R(s, sx, st_y + i*ROW_H - 0.04, sw, 0.010, BORDER)
        if i == 6:
            R(s, sx, st_y + i*ROW_H - 0.04, sw, 0.010, BORDER)

        y       = st_y + i * ROW_H
        bg      = PANEL if i % 2 == 0 else WHITE
        R(s, sx, y, sw, ROW_H - 0.02, bg)

        val_raw = row.get(key)
        suf     = "%" if key in ("efg_pct","usg_pct","three_pct") else ""
        lbl     = STAT_META[key][0]
        if pd.isna(val_raw):
            val_str = f"{0:.1f}{suf}"
            pct     = 0.0
            eff_pct = 0.0
            clr     = RED
        else:
            pct     = get_pct(val_raw, pct_arrays.get(key))
            eff_pct = pct if hib else (100 - pct if pct is not None else None)
            clr     = val_color(pct, hib)
            val_str = fmt(val_raw, 1, suf)

        mid_y = y + (ROW_H - 0.02)/2 - 0.11

        T(s, lbl,     sx+0.10,    mid_y, sw*0.28, 0.22, sz=8, clr=TEXT_MED)
        T(s, val_str, sx+sw*0.30, mid_y, sw*0.16, 0.22, sz=10, bold=True, clr=clr)

        bar_x = sx + sw*0.48
        bar_w = sw*0.32
        pct_bar(s, pct, bar_x, y + (ROW_H-0.02)/2 - 0.028, bar_w, 0.055, hib)

        if eff_pct is not None:
            T(s, f"{ordinal(int(eff_pct))} pct.",
              sx + sw*0.81, mid_y, sw*0.18, 0.22,
              sz=7.5, clr=clr, align=PP_ALIGN.RIGHT)

    # ── COL 2: Player Rating ──────────────────────────────────────────────────
    rx = C2L + PAD
    rw = C2W - PAD*2
    ry = BY + PAD

    T(s, "PLAYER RATING", rx, ry, rw, 0.28, sz=9, bold=True, clr=SCARLET)
    T(s, "Career production  ·  vs all D1",
      rx, ry+0.28, rw, 0.20, sz=7, clr=TEXT_LIGHT, italic=True)

    obpr_raw = row.get("obpr")
    dbpr_raw = row.get("dbpr")
    obpr_nan = pd.isna(obpr_raw)
    dbpr_nan = pd.isna(dbpr_raw)
    obpr   = 0.0 if obpr_nan else float(obpr_raw)
    dbpr   = 0.0 if dbpr_nan else float(dbpr_raw)
    bpr    = obpr + dbpr
    bpr_is_nan = obpr_nan or dbpr_nan

    card_w = rw
    card_h = 3.25
    card_y = ry + 0.54
    cx     = rx

    R(s, cx, card_y, card_w, card_h, PANEL, BORDER)
    R(s, cx, card_y, card_w, 0.44, SCARLET)
    T(s, "Career BPR", cx+0.18, card_y+0.06, card_w-0.30, 0.20,
      sz=8, bold=True, clr=WHITE)
    T(s, "Bayesian Performance Rating  ·  career production",
      cx+0.18, card_y+0.22, card_w-0.30, 0.16,
      sz=6.5, clr=WHITE, italic=True)

    for j, (stat_lbl, val) in enumerate([("Offensive", obpr), ("Defensive", dbpr)]):
        row_y = card_y + 0.44 + 0.14 + j * 0.50
        bg_r  = WHITE if j == 0 else PANEL
        R(s, cx+0.12, row_y, card_w-0.24, 0.46, bg_r)
        T(s, stat_lbl, cx+0.20, row_y+0.12, card_w*0.42, 0.24,
          sz=8.5, clr=TEXT_LIGHT)
        T(s, fmt(val,2), cx+card_w*0.44, row_y+0.10, card_w*0.50, 0.28,
          sz=13, bold=True, clr=TEXT_MED, align=PP_ALIGN.RIGHT)

    ov_y = card_y + 0.44 + 0.14 + 2*0.50 + 0.14
    R(s, cx+0.12, ov_y, card_w-0.24, 0.010, BORDER)
    ov_y += 0.16
    T(s, "Overall BPR", cx+0.14, ov_y, card_w-0.24, 0.24,
      sz=8, clr=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    ov_y += 0.32
    T(s, fmt(bpr, 2), cx+0.10, ov_y, card_w-0.20, 0.54,
      sz=30, bold=True, clr=TEXT_MED, align=PP_ALIGN.CENTER)
    ov_y += 0.60

    bpr_pct = 0.0 if bpr_is_nan else get_pct(bpr, pct_arrays.get("bpr"))
    pct_bar(s, bpr_pct, cx+0.18, ov_y, card_w-0.36, 0.055, True)
    if bpr_pct is not None:
        T(s, f"{ordinal(int(bpr_pct))} percentile",
          cx, ov_y+0.10, card_w, 0.22,
          sz=7.5, clr=val_color(bpr_pct), align=PP_ALIGN.CENTER, bold=True)

    # ── Scout Context ─────────────────────────────────────────────────────────
    sc_y = card_y + card_h + 0.18
    R(s, rx, sc_y, rw, 0.010, BORDER)
    sc_y += 0.14

    T(s, "CONTEXT", rx, sc_y, rw, 0.26, sz=8.5, bold=True, clr=SCARLET)
    ctx_rows = generate_scout_context(row)
    sc_y += 0.30
    LBL_W = 1.25
    for lbl, val, clr in ctx_rows:
        T(s, lbl, rx, sc_y, LBL_W, 0.38, sz=7.5, clr=TEXT_LIGHT)
        T(s, val, rx + LBL_W, sc_y, rw - LBL_W, 0.38, sz=8.5, bold=True, clr=clr)
        sc_y += 0.38

    # Footer
    R(s, 0, 7.38, 13.33, 0.12, PANEL)
    T(s, "2026 Portal Production Report",
      0.18, 7.39, 8, 0.11, sz=6.5, clr=TEXT_LIGHT)
    T(s, f"{group_name} #{rank} of {n_total}",
      11.0, 7.39, 2.15, 0.11, sz=6.5, clr=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


# ── Position-group section header ─────────────────────────────────────────────

def slide_group_header(prs, group_name, pool):
    """Section divider slide for a position group (Guards / Wings / Bigs)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    R(s, 0, 0, 13.33, 7.5, WHITE)

    LEFT_W  = 5.50
    RIGHT_W = 7.83

    # Left panel — crimson
    R(s, 0, 0, LEFT_W, 7.5, SCARLET)

    T(s, f"SECTION  {GROUP_ROMAN[group_name]}", 0.40, 0.50, LEFT_W - 0.55, 0.30,
      sz=10, bold=True, clr=RGBColor(0xFF, 0xAA, 0xAA))

    T(s, "TOP 15",             0.40, 1.00, LEFT_W - 0.55, 0.88, sz=54, bold=True, clr=WHITE)
    T(s, group_name.upper(),   0.40, 1.86, LEFT_W - 0.55, 0.88, sz=54, bold=True, clr=WHITE)

    R(s, 0.40, 2.86, 2.60, 0.040, WHITE)

    T(s, GROUP_ROLES_DESC[group_name],
      0.40, 2.98, LEFT_W - 0.55, 0.34, sz=12, clr=WHITE)
    T(s, "This past year's portal class  ·  ranked by career production (BPR)",
      0.40, 3.38, LEFT_W - 0.55, 0.30, sz=10, clr=RGBColor(0xFF, 0xCC, 0xCC), italic=True)

    # Right panel — white
    R(s, LEFT_W, 0, RIGHT_W, 7.5, WHITE)

    RX = LEFT_W + 0.42
    RW = RIGHT_W - 0.55

    T(s, f"SECTION {GROUP_ROMAN[group_name]}  ·  {group_name.upper()}",
      RX, 0.42, RW, 0.28, sz=8.5, bold=True, clr=TEXT_LIGHT)

    R(s, RX, 0.76, RW, 0.010, BORDER)

    T(s, "Coming up:", RX, 0.90, RW, 0.28, sz=13, bold=True, clr=SCARLET)

    list_y = 1.28
    for i, (_, p_row) in enumerate(pool.iterrows()):
        if i >= TOP_N_PER_GROUP:
            break
        name   = str(p_row.get("espn_name", ""))
        school = str(p_row.get("espn_team", "") or "")

        R(s, RX, list_y + 0.02, 0.26, 0.26, SCARLET_BG)
        num_box = T(s, str(i + 1), RX, list_y + 0.02, 0.26, 0.26,
          sz=8, bold=True, clr=SCARLET, align=PP_ALIGN.CENTER)
        num_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        T(s, name, RX + 0.34, list_y + 0.05, 3.60, 0.28,
          sz=10.5, bold=True, clr=TEXT)
        T(s, school, RX + 4.00, list_y + 0.07, RW - 4.05, 0.26,
          sz=9, clr=TEXT_LIGHT, italic=True)

        list_y += 0.385

    # Bottom strip
    R(s, 0, 7.15, 13.33, 0.35, RGBColor(0x1E, 0x1E, 0x1E))
    T(s, "2026 TRANSFER PORTAL PRODUCTION REPORT",
      0.30, 7.19, 13.0, 0.26, sz=7.5, bold=True,
      clr=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    T(s, f"SECTION  {GROUP_ROMAN[group_name]}", 12.30, 7.19, 0.95, 0.22,
      sz=7, bold=True, clr=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if not os.path.exists(MASTER_CSV):
        print(f"[ERROR] {MASTER_CSV} not found.")
        sys.exit(1)

    print("[DECK] Loading master CSV...")
    master = pd.read_csv(MASTER_CSV, dtype={"espn_id": str})
    print(f"  {len(master)} players loaded")

    games_n = pd.to_numeric(master["games"], errors="coerce")

    # Candidate pool = every in-portal player in the dataset (no BPR
    # requirement here — that's only applied later, at ranking time).
    pool_all = master[
        master["in_portal"].eq(True) &
        games_n.le(MAX_GAMES)
    ].copy()

    if EXCLUDE_FROM_TEAMS:
        excl_teams = {t.lower().strip() for t in EXCLUDE_FROM_TEAMS}
        excluded = pool_all[pool_all["espn_team"].str.lower().str.strip().isin(excl_teams)]
        if len(excluded):
            print(f"  Excluded (from {', '.join(sorted(EXCLUDE_FROM_TEAMS))}): "
                  f"{', '.join(sorted(excluded['espn_name']))}")
        pool_all = pool_all[~pool_all["espn_team"].str.lower().str.strip().isin(excl_teams)]

    if EXCLUDE_PLAYERS:
        excl = {n.lower().strip() for n in EXCLUDE_PLAYERS}
        pool_all = pool_all[
            ~pool_all["espn_name"].str.lower().str.strip().isin(excl)
        ]
        print(f"  Excluded:                    {', '.join(sorted(EXCLUDE_PLAYERS))}")

    print(f"  Portal pool (all in-portal players): {len(pool_all)}")

    print("\n[PCT] Building D1 percentile arrays...")
    pct_arrays = build_pct_arrays(master)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("[DECK] Building slides...")
    slide_title(prs, TOP_N_PER_GROUP)
    slide_explainer_structure(prs)
    slide_explainer_bpr(prs)

    n_players = 0
    for group_name, roles in POSITION_GROUPS:
        # No BPR → no rank: a player can be in the candidate pool without
        # BPR, but only BPR'd players are eligible to actually fill a
        # ranked Top-15 slot.
        grp = (pool_all[pool_all["role"].isin(roles) & pool_all["bpr"].notna()]
               .sort_values("bpr", ascending=False)
               .head(TOP_N_PER_GROUP)
               .reset_index(drop=True))
        n_players += len(grp)

        print(f"\n[DECK] Section {GROUP_ROMAN[group_name]} — {group_name} ({len(grp)} players)...")
        slide_group_header(prs, group_name, grp)
        slide_group_table(prs, group_name, grp, pct_arrays)

        for rank, (_, row) in enumerate(grp.iterrows(), start=1):
            slide_player(prs, rank, row, pct_arrays, len(grp), group_name)
            to_team = row.get("bt_to_team")
            status  = str(to_team) if pd.notna(to_team) else "Available"
            print(f"  {rank:2d}. {row['espn_name']:<28} "
                  f"role={str(row.get('role','')):<12}  "
                  f"bpr={float(row.get('bpr') or 0):.2f}  "
                  f"status={status}")

    out = os.path.join(OUTPUT, "iu_portal_report_by_position.pptx")
    prs.save(out)
    n_slides = 3 + len(POSITION_GROUPS) * 2 + n_players
    print(f"\n✅  Saved → {out}  ({n_slides} slides)")


if __name__ == "__main__":
    main()
